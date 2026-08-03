"""Build a 10-slide PM deck for the Category Discovery Engine project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Category_Discovery_Engine_Deck.pptx"

# Brand (Blinkit-adjacent green, not purple AI defaults)
GREEN = RGBColor(0x0D, 0x9F, 0x6E)
GREEN_DARK = RGBColor(0x05, 0x66, 0x4A)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF9, 0xF8)
CARD = RGBColor(0xEC, 0xF4, 0xF0)
AMBER = RGBColor(0xB5, 0x47, 0x08)
RED = RGBColor(0xB4, 0x23, 0x18)

W = Inches(13.333)
H = Inches(7.5)


def _set_run(run, *, size=14, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, color=CREAM):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def add_footer(slide, n: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(10), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Category Discovery Engine  ·  Why don't Blinkit users explore new categories?"
    _set_run(run, size=10, color=MUTED)
    num = slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3))
    p2 = num.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"{n} / 10"
    _set_run(run2, size=10, color=MUTED)


def title_block(slide, kicker: str, title: str, subtitle: str | None = None):
    k = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.35))
    p = k.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = kicker.upper()
    _set_run(run, size=12, bold=True, color=GREEN)

    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.65), Inches(12.2), Inches(0.7))
    p = t.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=28, bold=True, color=INK)

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.4))
        p = s.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        _set_run(run, size=14, color=MUTED)


def card(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.08
    return shape


def text_in(slide, left, top, width, height, lines: list[tuple[str, dict]]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, style in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(style.get("after", 4))
        run = p.add_run()
        run.text = text
        _set_run(
            run,
            size=style.get("size", 13),
            bold=style.get("bold", False),
            color=style.get("color", INK),
        )
    return box


def bullet_block(slide, left, top, width, height, bullets: list[str], size=13):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {b}"
        _set_run(run, size=size, color=INK)
    return box


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ------------------------------------------------------------------
    # SLIDE 1 — CONTEXT · THE PLATFORM TODAY
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "01  ·  Context",
        "The platform today",
        "Blinkit wins on speed for groceries — but users rarely expand into new categories.",
    )
    card(s, Inches(0.55), Inches(1.9), Inches(4.0), Inches(4.5))
    text_in(
        s,
        Inches(0.75),
        Inches(2.1),
        Inches(3.6),
        Inches(4.1),
        [
            ("What Blinkit is great at", {"size": 15, "bold": True, "color": GREEN_DARK, "after": 10}),
            ("•  10–20 min delivery for staples", {"size": 13, "after": 6}),
            ("•  Habitual reorder / grocery mental model", {"size": 13, "after": 6}),
            ("•  Dense SKU coverage in core categories", {"size": 13, "after": 6}),
            ("•  Strong late-night convenience use case", {"size": 13, "after": 14}),
            ("Primary research question", {"size": 15, "bold": True, "color": GREEN_DARK, "after": 8}),
            ("Why don't Blinkit users explore new categories?", {"size": 14, "bold": True, "after": 6}),
        ],
    )
    card(s, Inches(4.75), Inches(1.9), Inches(8.0), Inches(4.5), WHITE)
    text_in(
        s,
        Inches(5.0),
        Inches(2.1),
        Inches(7.5),
        Inches(4.1),
        [
            ("The growth tension", {"size": 15, "bold": True, "color": GREEN_DARK, "after": 10}),
            ("Users treat Blinkit as a narrow utility (milk, bread, snacks) — not a multi-category marketplace.", {"size": 14, "after": 12}),
            ("Business impact", {"size": 14, "bold": True, "after": 6}),
            ("•  Caps basket size & category mix per MAC", {"size": 13, "after": 4}),
            ("•  Under-monetizes Home, Beauty, Electronics, Pet, Pharma-adjacent", {"size": 13, "after": 4}),
            ("•  Recommendations/promos rarely address why users refuse unfamiliar categories", {"size": 13, "after": 12}),
            ("Evidence gap today", {"size": 14, "bold": True, "after": 6}),
            ("Feedback is scattered across Play/App Store, Reddit, social, forums — product teams lack a systematic barrier view.", {"size": 13, "after": 4}),
        ],
    )
    add_footer(s, 1)

    # ------------------------------------------------------------------
    # SLIDE 2 — MARKET LANDSCAPE
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "02  ·  Market landscape",
        "Competition & sentiment",
        "Quick-commerce is a race for wallet share — category expansion is the next battleground.",
    )
    # three columns
    cols = [
        (
            "Competitive set",
            [
                "Zepto / Instamart / Swiggy — similar speed narrative",
                "Kirana + supermarket still win trust for non-grocery",
                "Amazon / Flipkart for electronics & branded home",
                "Category expansion = differentiation beyond ETA",
            ],
        ),
        (
            "Sentiment themes (corpus)",
            [
                "Quality distrust dominates non-grocery talk (211 mentions)",
                "Price / fee uncertainty (133)",
                "Coverage & ETA anxiety (131)",
                "Assortment gaps teach 'not for this' (67)",
                "Praise is often generic ('good app') — low insight density",
            ],
        ),
        (
            "Implication",
            [
                "Speed alone won't unlock Home / Beauty / Electronics",
                "Trust, discovery adjacency, and price clarity are the wedges",
                "Snacks is the easiest expansion adjacency to grocery",
                "Home needs risk reducers (returns / quality cues)",
            ],
        ),
    ]
    for i, (head, bullets) in enumerate(cols):
        left = Inches(0.55 + i * 4.15)
        card(s, left, Inches(1.9), Inches(3.95), Inches(4.5), WHITE if i != 1 else CARD)
        text_in(
            s,
            left + Inches(0.2),
            Inches(2.1),
            Inches(3.55),
            Inches(0.4),
            [(head, {"size": 15, "bold": True, "color": GREEN_DARK})],
        )
        bullet_block(s, left + Inches(0.2), Inches(2.55), Inches(3.55), Inches(3.6), bullets, size=12)
    add_footer(s, 2)

    # ------------------------------------------------------------------
    # SLIDE 3 — SEGMENTATION · PERSONAS · HYPOTHESES
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "03  ·  Segmentation · Personas · Hypotheses",
        "Who stays stuck — and what the evidence says",
        "Corpus: 2,818 reviews · 643 exploration-relevant · segments from KMeans on feedback.",
    )
    # personas strip
    personas = [
        ("Routine Buyers", "554", "Same milk/bread loop; never browse"),
        ("Price Sensitive", "1,736", "Fee/markup fear blocks trial"),
        ("Explorers", "348", "Open to try if trust + clarity"),
        ("Impulse Buyers", "180", "Attach if surfaced in-session"),
    ]
    for i, (name, n, blurb) in enumerate(personas):
        left = Inches(0.55 + i * 3.15)
        card(s, left, Inches(1.85), Inches(3.0), Inches(1.55), WHITE)
        text_in(
            s,
            left + Inches(0.15),
            Inches(1.95),
            Inches(2.7),
            Inches(1.3),
            [
                (name, {"size": 13, "bold": True, "color": GREEN_DARK, "after": 2}),
                (f"n≈{n}", {"size": 11, "color": MUTED, "after": 4}),
                (blurb, {"size": 11, "after": 2}),
            ],
        )
    # hypotheses table-like
    card(s, Inches(0.55), Inches(3.55), Inches(12.2), Inches(3.0), WHITE)
    hyps = [
        ("H1 Discovery", "8", "Adjacent categories not surfaced in grocery sessions"),
        ("H2 Quality trust", "211", "Non-grocery quality/authenticity feels riskier"),
        ("H3 Price clarity", "133", "All-in fees unclear vs offline"),
        ("H4 Habit reorder", "36", "Power reorderers never leave the loop"),
        ("H5 Assortment", "67", "Missing SKUs teach 'Blinkit isn't for this'"),
        ("H6 Coverage/ETA", "131", "ETA anxiety reduces unfamiliar-category trial"),
    ]
    text_in(
        s,
        Inches(0.75),
        Inches(3.7),
        Inches(11.8),
        Inches(0.35),
        [("Hypotheses × evidence mentions (barrier-linked)", {"size": 14, "bold": True, "color": GREEN_DARK})],
    )
    for i, (hid, mentions, stmt) in enumerate(hyps):
        col = i % 3
        row = i // 3
        left = Inches(0.75 + col * 4.05)
        top = Inches(4.15 + row * 1.1)
        text_in(
            s,
            left,
            top,
            Inches(3.9),
            Inches(1.0),
            [
                (f"{hid}  ·  {mentions} mentions", {"size": 12, "bold": True, "color": AMBER if int(mentions) >= 100 else INK, "after": 2}),
                (stmt, {"size": 11, "color": MUTED}),
            ],
        )
    add_footer(s, 3)

    # ------------------------------------------------------------------
    # SLIDE 4 — PROBLEM CANVAS
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "04  ·  Problem canvas", "Problem canvas", "From jobs-to-be-done to unmet needs — grounded in exploration tags.")
    blocks = [
        ("Customer job", "When I need something outside my usual list, I want to trust assortment & quality so I can order without second-guessing."),
        ("Situation", "Short grocery reorder sessions under time pressure; mental model = 'Blinkit = staples'."),
        ("Desired outcome", "Serendipitous but relevant discovery + bounded risk (returns, fees, quality cues)."),
        ("Current workaround", "Order only staples on Blinkit; buy Home/Beauty/Electronics elsewhere."),
        ("Pains / barriers", "Quality distrust · price uncertainty · weak in-session discovery · assortment gaps · ETA anxiety."),
        ("Unmet needs (P0)", "1) Trusted paths into non-grocery during grocery sessions  2) Risk reducers for first buys outside grocery."),
    ]
    for i, (h, body) in enumerate(blocks):
        col = i % 3
        row = i // 3
        left = Inches(0.55 + col * 4.15)
        top = Inches(1.9 + row * 2.4)
        card(s, left, top, Inches(3.95), Inches(2.2), WHITE)
        text_in(
            s,
            left + Inches(0.2),
            top + Inches(0.2),
            Inches(3.55),
            Inches(1.8),
            [
                (h, {"size": 13, "bold": True, "color": GREEN_DARK, "after": 8}),
                (body, {"size": 12, "after": 2}),
            ],
        )
    add_footer(s, 4)

    # ------------------------------------------------------------------
    # SLIDE 5 — SOLUTION IDEATION · RICE
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "05  ·  Solution ideation",
        "Ideation & RICE prioritisation",
        "Four experiment seeds ranked; two MVPs selected for Prototype Lab.",
    )
    headers = ["Idea", "R", "I", "C", "E", "RICE", "Decision"]
    rows = [
        ["Snacks discovery rail (after grocery add)", "9", "8", "3", "2", "108", "P0 MVP"],
        ["Home first-buy quality guarantee", "8", "8", "4", "3", "85", "P0 MVP"],
        ["All-in price clarity on new-category PDPs", "7", "6", "3", "2", "63", "P1"],
        ["Routine-break prompt for power reorderers", "6", "5", "4", "3", "40", "P1"],
    ]
    # header row
    for i, h in enumerate(headers):
        widths = [5.2, 0.7, 0.7, 0.7, 0.7, 1.0, 1.6]
        left = Inches(0.55 + sum(widths[:i]))
        card(s, left, Inches(1.9), Inches(widths[i] - 0.08), Inches(0.45), GREEN_DARK)
        text_in(
            s,
            left + Inches(0.08),
            Inches(1.97),
            Inches(widths[i] - 0.2),
            Inches(0.35),
            [(h, {"size": 12, "bold": True, "color": WHITE})],
        )
    for r, row in enumerate(rows):
        top = Inches(2.45 + r * 0.7)
        fill = CARD if r < 2 else WHITE
        for i, cell in enumerate(row):
            widths = [5.2, 0.7, 0.7, 0.7, 0.7, 1.0, 1.6]
            left = Inches(0.55 + sum(widths[:i]))
            card(s, left, top, Inches(widths[i] - 0.08), Inches(0.62), fill)
            text_in(
                s,
                left + Inches(0.08),
                top + Inches(0.15),
                Inches(widths[i] - 0.2),
                Inches(0.4),
                [(cell, {"size": 12, "bold": i == 0 or i == 6, "color": GREEN_DARK if i == 6 and r < 2 else INK})],
            )
    text_in(
        s,
        Inches(0.55),
        Inches(5.4),
        Inches(12.2),
        Inches(1.2),
        [
            ("RICE note", {"size": 13, "bold": True, "color": GREEN_DARK, "after": 6}),
            ("Reach = grocery-session volume that can see the intervention · Impact = attach / first-category conversion · Confidence from barrier evidence (H2 strongest) · Effort = eng + design + ops. Snacks rail wins on adjacency + low trust friction; Home guarantee attacks the #1 barrier (quality distrust).", {"size": 12, "after": 4}),
        ],
    )
    add_footer(s, 5)

    # ------------------------------------------------------------------
    # SLIDE 6 — SOLUTION FLOW · AI WEDGE
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "06  ·  The solution",
        "Flow, mechanism & the AI wedge",
        "User workflow (product) + insight engine (how we know what to build).",
    )
    # user flow boxes
    flow = [
        ("1. Grocery add", "User reorders staples"),
        ("2. Discovery rail", "'Also useful tonight' snacks"),
        ("3. Trust cues", "Ratings · all-in price · badge"),
        ("4. Attach / buy", "Non-grocery SKU in-session"),
        ("5. First-buy path", "Home + easy return guarantee"),
    ]
    for i, (h, b) in enumerate(flow):
        left = Inches(0.4 + i * 2.55)
        card(s, left, Inches(1.9), Inches(2.4), Inches(1.5), WHITE)
        text_in(
            s,
            left + Inches(0.12),
            Inches(2.05),
            Inches(2.15),
            Inches(1.2),
            [
                (h, {"size": 12, "bold": True, "color": GREEN_DARK, "after": 4}),
                (b, {"size": 11, "color": MUTED}),
            ],
        )
        if i < 4:
            arrow = s.shapes.add_textbox(left + Inches(2.25), Inches(2.4), Inches(0.35), Inches(0.4))
            p = arrow.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = ">"
            _set_run(run, size=16, bold=True, color=GREEN)

    card(s, Inches(0.55), Inches(3.65), Inches(12.2), Inches(2.9), CARD)
    text_in(
        s,
        Inches(0.75),
        Inches(3.8),
        Inches(11.8),
        Inches(2.6),
        [
            ("AI wedge — Category Discovery Engine", {"size": 15, "bold": True, "color": GREEN_DARK, "after": 8}),
            ("Collect multi-source feedback  →  clean & embed  →  themes / sentiment / segments  →  exploration tags (barriers + signals)  →  synthesis (JTBD, hypotheses, experiments)  →  Streamlit research console + Prototype Lab", {"size": 13, "after": 10}),
            ("Mechanism: rule-based exploration tagger (explainable) + BERTopic themes + optional LLM polish. Output is decision-ready: ranked barriers, category opportunities, and linked MVP experiments — not another generic dashboard.", {"size": 13, "after": 8}),
            ("Human-in-the-loop: Validation Desk approve/reject hypotheses before treating claims as decision-grade.", {"size": 13, "after": 4}),
        ],
    )
    add_footer(s, 6)

    # ------------------------------------------------------------------
    # SLIDE 7 — SYSTEM ARCHITECTURE · MVP BUILD
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "07  ·  System architecture",
        "MVP build — 13-stage pipeline + Streamlit console",
        "Primary UX: Streamlit. Optional: Theme Explorer API (advanced).",
    )
    stages = [
        ("Collect", "Play · App Store\nReddit · YouTube"),
        ("Merge &\nclean", "Soft-dedupe →\npreprocessed CSV"),
        ("NLP", "Embed · themes\nsentiment · segments"),
        ("Explore\n& synthesize", "Tags → barriers\nJTBD · experiments"),
        ("Serve", "Findings · Validation\nPrototype Lab"),
    ]
    for i, (h, b) in enumerate(stages):
        left = Inches(0.5 + i * 2.55)
        card(s, left, Inches(1.95), Inches(2.4), Inches(2.0), WHITE)
        text_in(
            s,
            left + Inches(0.15),
            Inches(2.1),
            Inches(2.1),
            Inches(1.7),
            [
                (h.replace("\n", " "), {"size": 13, "bold": True, "color": GREEN_DARK, "after": 6}),
                (b.replace("\n", " · "), {"size": 11, "color": MUTED}),
            ],
        )
    card(s, Inches(0.55), Inches(4.2), Inches(6.0), Inches(2.4), CARD)
    bullet_block(
        s,
        Inches(0.75),
        Inches(4.35),
        Inches(5.6),
        Inches(2.1),
        [
            "Artifacts: merged → themes / sentiment / segments / exploration_tags → synthesis.json",
            "Parity: NLP + exploration aligned on preprocessed corpus (2,818)",
            "Gold set + barrier F1 eval (heuristic ≈ 0.88 micro-F1)",
            "CI smoke on PRs · human review checklist on Validation Desk",
        ],
        size=12,
    )
    card(s, Inches(6.75), Inches(4.2), Inches(6.0), Inches(2.4), WHITE)
    bullet_block(
        s,
        Inches(6.95),
        Inches(4.35),
        Inches(5.6),
        Inches(2.1),
        [
            "Streamlit IA: Findings · Opportunities · Validation · Live Pipeline · Try-it · Prototype Lab · Evidence · Admin",
            "MVP mocks: Snacks rail + Home first-buy guarantee",
            "Exports: experiment brief Markdown · stakeholder event log",
            "Run: python main.py --skip-collect  →  streamlit run app.py",
        ],
        size=12,
    )
    add_footer(s, 7)

    # ------------------------------------------------------------------
    # SLIDE 8 — MVP METRICS
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "08  ·  MVP metrics",
        "North star, leading, lagging & guardrails",
        "Instrument before scaling either MVP beyond Prototype Lab.",
    )
    metrics = [
        ("North star", GREEN_DARK, WHITE, "% of MACs who purchase from ≥1 new product category each month"),
        ("Leading — Snacks rail", GREEN, WHITE, "Attach rate of non-grocery SKU in same grocery session · rail CTR · PDP→ATC on rail SKUs"),
        ("Leading — Home guarantee", GREEN, WHITE, "First-time Home conversion among grocery-only cohorts · guarantee view→buy"),
        ("Lagging", AMBER, WHITE, "Category mix per MAC · repeat rate in new category · 30-day expanded basket value"),
        ("Guardrails", RED, WHITE, "Checkout completion drop ≤2% (rail) · return-rate spike <3pp (guarantee) · no surge in 'hidden fees' tickets"),
        ("Insight quality", MUTED, WHITE, "Barrier micro-F1 on gold · hypothesis triangulation strength · % hypotheses human-approved"),
    ]
    for i, (h, bg, fg, body) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left = Inches(0.55 + col * 4.15)
        top = Inches(1.9 + row * 2.4)
        card(s, left, top, Inches(3.95), Inches(2.2), WHITE)
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.95), Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = bg
        hdr.line.fill.background()
        hdr.adjustments[0] = 0.08
        text_in(
            s,
            left + Inches(0.15),
            top + Inches(0.08),
            Inches(3.65),
            Inches(0.35),
            [(h, {"size": 13, "bold": True, "color": fg})],
        )
        text_in(
            s,
            left + Inches(0.2),
            top + Inches(0.65),
            Inches(3.55),
            Inches(1.4),
            [(body, {"size": 13})],
        )
    add_footer(s, 8)

    # ------------------------------------------------------------------
    # SLIDE 9 — PROTOTYPE LAB (user didn't specify; natural continuation)
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "09  ·  Prototype Lab",
        "Two clickable MVPs for stakeholder demos",
        "Research prototypes — not production Blinkit. Linked to synthesis experiments.",
    )
    card(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(4.5), WHITE)
    text_in(
        s,
        Inches(0.75),
        Inches(2.1),
        Inches(5.6),
        Inches(4.1),
        [
            ("MVP 1 — Grocery → Snacks rail", {"size": 16, "bold": True, "color": GREEN_DARK, "after": 8}),
            ("Experiment: exp_discover_rail  ·  Barrier: hard to discover", {"size": 12, "color": MUTED, "after": 8}),
            ("After staples land in cart, surface 'Also useful tonight' snacks with freshness badge + all-in price + ratings trust panel.", {"size": 13, "after": 8}),
            ("Primary metric: attach rate of non-grocery SKU in-session", {"size": 13, "bold": True, "after": 4}),
            ("Guardrail: checkout completion must not drop >2%", {"size": 12, "after": 10}),
            ("Why snacks first: highest adjacency to grocery; lowest quality-fear vs Home/Electronics.", {"size": 12, "color": MUTED}),
        ],
    )
    card(s, Inches(6.75), Inches(1.9), Inches(6.0), Inches(4.5), CARD)
    text_in(
        s,
        Inches(6.95),
        Inches(2.1),
        Inches(5.6),
        Inches(4.1),
        [
            ("MVP 2 — Home first-buy guarantee", {"size": 16, "bold": True, "color": GREEN_DARK, "after": 8}),
            ("Experiment: exp_first_buy_guarantee  ·  Barrier: quality distrust (#1)", {"size": 12, "color": MUTED, "after": 8}),
            ("First Home purchase shows quality score + one-tap easy return — bounds the risk that blocks trial.", {"size": 13, "after": 8}),
            ("Primary metric: first-time category conversion (grocery-only cohort)", {"size": 13, "bold": True, "after": 4}),
            ("Guardrail: return-rate spike <3pp vs control", {"size": 12, "after": 10}),
            ("Why Home: high blocked intent + quality fear is the dominant barrier in the corpus.", {"size": 12, "color": MUTED}),
        ],
    )
    add_footer(s, 9)

    # ------------------------------------------------------------------
    # SLIDE 10 — ASK / NEXT STEPS
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(
        s,
        "10  ·  Ask & next steps",
        "What we need to move from research console → shipped experiment",
        "Decision-grade path: approve hypotheses → instrument MVPs → 2-week city test.",
    )
    card(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(4.5), WHITE)
    text_in(
        s,
        Inches(0.75),
        Inches(2.1),
        Inches(5.6),
        Inches(4.1),
        [
            ("The ask", {"size": 16, "bold": True, "color": GREEN_DARK, "after": 10}),
            ("1. Approve H2 + snacks/home opportunity framing on Validation Desk", {"size": 13, "after": 6}),
            ("2. Green-light 2-week A/B (or city holdout) for Snacks rail", {"size": 13, "after": 6}),
            ("3. Parallel design spike for Home first-buy guarantee", {"size": 13, "after": 6}),
            ("4. Assign owners: Discovery / Merch / Trust & Returns", {"size": 13, "after": 12}),
            ("Success in 30 days", {"size": 14, "bold": True, "color": GREEN_DARK, "after": 6}),
            ("Measurable attach lift on rail + no guardrail breach; decision to scale or kill.", {"size": 13}),
        ],
    )
    card(s, Inches(6.75), Inches(1.9), Inches(6.0), Inches(4.5), CARD)
    text_in(
        s,
        Inches(6.95),
        Inches(2.1),
        Inches(5.6),
        Inches(4.1),
        [
            ("Near-term roadmap", {"size": 16, "bold": True, "color": GREEN_DARK, "after": 10}),
            ("Week 1–2: Instrument metrics + finalize experiment brief", {"size": 13, "after": 6}),
            ("Week 3–4: Snacks rail experiment live in 1–2 cities", {"size": 13, "after": 6}),
            ("Week 5–6: Readout · iterate copy/trust cues · queue Home guarantee", {"size": 13, "after": 6}),
            ("Ongoing: Refresh corpus via main.py · keep synthesis + human review current", {"size": 13, "after": 12}),
            ("Artifact links", {"size": 14, "bold": True, "color": GREEN_DARK, "after": 6}),
            ("streamlit run app.py  ·  output/synthesis.json  ·  docs/architecture.md", {"size": 12, "color": MUTED}),
        ],
    )
    add_footer(s, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
